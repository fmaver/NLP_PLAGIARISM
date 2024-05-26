"""
This script will be used to analyze the documents and extract the text from them.
"""

import os

import docx
from pptx import Presentation
from PyPDF2 import PdfReader

from plagiarism.adapters import file_handler


class DocumentAnalyzer:
    def __init__(self):
        self.docx_reader = DocxReader()
        self.pdf_reader = PDFReader()
        self.pptx_reader = PPTXReader()
        self.file_handler = file_handler.FileHandler()

    def txt_file_exists(self, output_path, file_name: str) -> bool:
        """
        This function checks if a text file already exists in the output folder.
        Args:
            file_name: Name of the file
            output_path: Path to the text file

        Returns: True if the file exists, False otherwise
        """
        return os.path.exists(os.path.join(output_path, file_name))

    def save_txt_file_into(self, output_folder: str, text: str, file_name: str, extension: str) -> str:
        """
        This function saves the text into a txt file.
        Args:
            file_name: Name of the file
            output_folder: folder to save the text file
            text: Text to be saved
            extension: File extension

        Returns: Path to the saved text file
        """
        output_file = os.path.join(output_folder, file_name.replace(extension, ".txt"))

        with open(output_file, "w", encoding="utf-8") as file:
            file.write(text)

        if self.file_handler.repeated_files(output_file, output_folder):
            os.remove(output_file)
            return output_file

        print(f"Text file {output_file} has been saved into {output_file} folder")
        return output_file

    def convert_files_to_txt(self, files: list[str], output_folder: str):
        """
        This function converts different types of documents (PDF, DOCX, PPTX) to text files
        and saves them in a 'dataset' folder.
        Args:
            files: List of files to be converted to text
            output_folder: Folder to save the text files

        Returns: None
        """
        for file in files:
            extension = ""
            file_name = os.path.basename(file)
            if file_name.endswith(".docx"):
                text = self.docx_reader.read_docx(file)
                extension = ".docx"
            elif file_name.endswith(".pdf"):
                text = self.pdf_reader.read_pdf(file)
                extension = ".pdf"
            elif file_name.endswith(".pptx"):
                text = self.pptx_reader.read_pptx(file)
                extension = ".pptx"
            elif file_name.endswith(".doc"):
                extension = ".doc"
                text = "This file format is not supported"

            else:
                raise ValueError(f"Unsupported file format: {file_name}")

            if not extension == ".doc":
                # check if the output folder don't exist
                if not (self.txt_file_exists(output_folder, file_name.replace(extension, ".txt"))):
                    self.save_txt_file_into(output_folder, text, file_name.replace(extension, ".txt"), extension)
                else:
                    output_path = os.path.join(output_folder, file_name.replace(extension, ".txt"))
                    print(f"Text file {output_path} already exists in {output_folder} folder")


class PDFReader:
    """
    PDFReader class that reads a pdf file and convert it into a txt file.
    Arguments for methods:
        pdf_path {str} -- Path to the pdf file.
        txt_path {str} -- Path to save the txt file.
    """

    def read_pdf(self, pdf_path) -> str:
        """
        Read the pdf file and convert it into a text file.
        Args:
            pdf_path: Path to the pdf file.
        Returns: Text extracted from the pdf file.
        """
        text = ""
        with open(pdf_path, "rb") as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        return text

    def get_pdf_name(self, pdf_path):
        """
        Get the name of the pdf file.
        Args:
            pdf_path: Path to the pdf file.
        Returns: Name of the pdf file.
        """
        return os.path.basename(pdf_path)


class PPTXReader:
    """
    PPTReader class that reads a pptx file and convert it into a txt file.
    Arguments for methods:
        ppt_path {str} -- Path to the pptx file.
        txt_path {str} -- Path to save the txt file.
    """

    def read_pptx(self, pptx_path) -> str:
        """
        Read the pptx file and convert it into a text file.
        Args:
            pptx_path: Path to the pptx file.
        Returns: Text extracted from the pptx file.
        """
        text = ""
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text

    def get_pptx_name(self, pptx_path):
        """
        Get the name of the pptx file.
        Args:
            pptx_path: Path to the pptx file.
        Returns: Name of the pptx file.
        """
        return os.path.basename(pptx_path)


class DocxReader:
    """
    DocxReader class that reads a docx file and convert it into a txt file.
    Arguments for methods:
        docx_path {str} -- Path to the docx file.
        txt_path {str} -- Path to save the txt file.
    """

    def read_docx(self, docx_path) -> str:
        """
        Read the docx file and convert it into a text file.
        Args:
            docx_path: Path to the docx file.
        Returns: Text extracted from the docx file.
        """
        doc = docx.Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"  # Append newline character after each paragraph
        return text

    def get_docx_name(self, docx_path):
        """
        Get the name of the docx file.
        Args:
            docx_path: Path to the docx file.
        Returns: Name of the docx file.
        """
        return os.path.basename(docx_path)
